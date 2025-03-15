import os
import fitz  # PyMuPDF
import docx
import pptx
import cloudinary.uploader
from rest_framework.response import Response
from rest_framework.views import APIView
from rest_framework.parsers import MultiPartParser, FormParser
from django.core.files.storage import default_storage
from django.conf import settings

class FileUploadAPIView(APIView):
    parser_classes = (MultiPartParser, FormParser)

    def post(self, request, *args, **kwargs):
        file = request.FILES.get('file')

        if not file:
            return Response({"error": "No file uploaded"}, status=400)

        file_ext = os.path.splitext(file.name)[1].lower()

        # L∆∞u file t·∫°m th·ªùi
        file_path = default_storage.save(file.name, file)
        full_file_path = default_storage.path(file_path)

        print(f"üìÇ File path: {full_file_path}")  # Debug
        try:
            if file_ext == ".pdf":
                text = self.read_pdf(file_path)
            elif file_ext == ".docx":
                text = self.read_docx(file_path)
            elif file_ext == ".pptx":
                text = self.read_pptx(file_path)
            else:
                return Response({"error": "Unsupported file type"}, status=400)
        finally:
            # X√≥a file sau khi ƒë·ªçc xong
            default_storage.delete(file_path)

        return Response({"filename": file.name, "content": text})

    def read_pdf(self, file_path):
        text = ""
        with fitz.open(file_path) as pdf:
            for page in pdf:
                text += page.get_text("text")
        return text

    def read_docx(self, file_path):
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])

    def read_pptx(self, file_path):
        ppt = pptx.Presentation(file_path)
        text = []
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
class UploadImageView(APIView):
    parser_classes = (MultiPartParser, FormParser)

    def post(self, request, *args, **kwargs):
        file = request.FILES.get('file')
        if file:
            upload_result = cloudinary.uploader.upload(file)
            return Response({'url': upload_result['secure_url']})
        return Response({'error': 'No file provided'}, status=400)